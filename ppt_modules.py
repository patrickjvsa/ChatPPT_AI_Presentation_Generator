"""
Este archivo contiene algunos modulos creados con el fin de trabajar con archivos pptx.
Fue programado teniendo en mente su utilización en aplicaciones de IA.

Author: Patrick Vásquez <pvasquezs@fen.uchile.cl>
Date: 30/10/2023
"""

class Presentation():
    
    def __init__(self, title, author, template = "template.pptx"):
        """
        Inicializa una nueva presentación.

        Args:
        - title (str): El título de la presentación.
        - author (str): El autor de la presentación.
        - template (str): La ruta al archivo de plantilla. Revisa dentro de la carpeta templates. Por defecto, "template.pptx".

        Returns:
        - None
        """
        from pptx import Presentation
        self.title = title
        self.author = author
        self.template = "./templates/" + template
        self.presentation = Presentation(self.template)

    def add_frontpage(self):
        """
        Añade una diapositiva de portada a la presentación con el título y autor especificados.

        Args:
        - titulo (str): El título de la presentación.
        - autor (str): El autor de la presentación.

        Returns:
        - None
        """
        from pptx import Presentation
        slide_layout = self.presentation.slide_layouts[0]  
        slide = self.presentation.slides.add_slide(slide_layout)
        content_placeholder = slide.placeholders[0]
        content_placeholder.text = f"{self.title}"
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = f"{self.author}"

    def add_index(self, indice):
        """
        Añade un índice a la presentación.

        Args:
        - indice (str): El índice a agregar.

        Returns:
        - None
        """
        from pptx import Presentation
        slide_layout = self.presentation.slide_layouts[1]  
        slide = self.presentation.slides.add_slide(slide_layout)
        content_placeholder = slide.placeholders[0]
        content_placeholder.text = "Indice"
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = f"{indice}"

    def save_presentation(self, name="my_presentation.pptx"):
        """
        Guarda la presentación actual en un archivo con el nombre especificado.
        
        Args:
        - name (str): El nombre del archivo de la presentación. Por defecto es "my_presentation.pptx".
        Returns:
        - None
        """
        from pptx import Presentation
        self.presentation.save("./output/" + name)
        # cuando exista la webapp, hay que mover ese decorador output del modulo al archivo app.py

    def create_photo_slide_cat(self):
        """
        Crea una diapositiva con una foto de un gato.

        Args:
        - None

        Returns:
        - None
        """
        from pptx import Presentation
        from pptx.util import Inches
        slide_layout = self.presentation.slide_layouts[3]  
        slide = self.presentation.slides.add_slide(slide_layout)
        content_placeholder = slide.placeholders[0]
        content_placeholder.text = "Foto de gato"
        # Agrega una imagen a la diapositiva
        left = Inches(2)  # Ajusta la posición izquierda de la imagen
        top = Inches(3)   # Ajusta la posición superior de la imagen
        width = Inches(4) # Ajusta el ancho de la imagen
        height = Inches(3) # Ajusta la altura de la imagen
        image_path = "./templates/cat.jpg"  # Ruta a la imagen que deseas agregar
        pic = slide.shapes.add_picture(image_path, left, top, width, height)

    def create_photo_and_text_slide(self, title, text, image_path = "./dataset/Image_1.jpg"):
        """
        Crea una diapositiva con una foto y un texto.

        Args:
        - title (str): El título de la diapositiva.
        - text (str): El texto a agregar a la diapositiva.
        - image_path (str): La ruta a la imagen que deseas agregar.

        Returns:
        - None
        """
        from pptx import Presentation
        from pptx.util import Inches
        slide_layout = self.presentation.slide_layouts[3]  
        slide = self.presentation.slides.add_slide(slide_layout)

        # añade el título
        content_placeholder = slide.placeholders[0]
        content_placeholder.text = title

        # agrega una imagen a la diapositiva
        left = Inches(2)  # Ajusta la posición izquierda de la imagen
        top = Inches(3)   # Ajusta la posición superior de la imagen
        width = Inches(4) # Ajusta el ancho de la imagen
        height = Inches(3) # Ajusta la altura de la imagen

        #agrega el texto
        content_placeholder = slide.placeholders[2]
        content_placeholder.text = text

        pic = slide.shapes.add_picture(image_path, left, top, width, height)

# Ejemplo de uso
if __name__ == "__main__":

    import text_modules as tm

    print("Ejecutando programa principal...")

    query = 'pinguinos emperadores'
    autor = 'Patrick Vásquez'

    llm = tm.OpenAICompleter()
    titulo = llm.complete_text("Genera un titulo para una presentación sobre " + query + ".")

    presentation = Presentation(titulo, autor)

    print("Plantilla cargada con éxito.")

    # presentation.add_frontpage()
    # print("Portada agregada con éxito.")
    # presentation.add_index("Índice")
    # print("Índice agregado con éxito.")
    # presentation.create_photo_slide_cat()
    # print("Gato agregado con éxito.")
    # presentation.save_presentation("my_presentation.pptx")
    # print("Presentación guardada con éxito.") 

